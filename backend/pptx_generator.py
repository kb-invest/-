from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os
from datetime import datetime

def generate_briefing_pptx(property_data: dict, output_path: str):
    """
    브리핑용 PPTX 제안서 생성
    화려한 디자인, 프레젠테이션용
    """
    # 템플릿 로드
    template_path = os.path.join(os.path.dirname(__file__), '..', 'template_briefing.pptx')
    
    if os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    
    # 기존 슬라이드 삭제 (템플릿 기반으로 새로 생성)
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    
    # 1. 표지 슬라이드
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃
    
    # 배경색 설정
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 35, 53)  # 다크 블루
    
    # 프로젝트명
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = property_data.get('project_name', '상업용 부동산 투자 기회')
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # 건물명
    building_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(0.8))
    building_frame = building_box.text_frame
    building_frame.text = property_data.get('building_name', '')
    building_para = building_frame.paragraphs[0]
    building_para.font.size = Pt(36)
    building_para.font.color.rgb = RGBColor(255, 255, 255)
    building_para.alignment = PP_ALIGN.CENTER
    
    # 매각 제안서
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.3), Inches(8), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "매각 제안서"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(200, 200, 200)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # 매각가
    price_box = slide.shapes.add_textbox(Inches(3), Inches(4.2), Inches(4), Inches(0.8))
    price_frame = price_box.text_frame
    price_text = f"예상 매각가\n{property_data.get('sale_price', '')}원"
    price_frame.text = price_text
    for para in price_frame.paragraphs:
        para.font.size = Pt(28)
        para.font.bold = True
        para.font.color.rgb = RGBColor(255, 215, 0)  # 골드
        para.alignment = PP_ALIGN.CENTER
    
    # 주소
    addr_box = slide.shapes.add_textbox(Inches(1), Inches(5.3), Inches(8), Inches(0.6))
    addr_frame = addr_box.text_frame
    addr_frame.text = property_data.get('address', '')
    addr_para = addr_frame.paragraphs[0]
    addr_para.font.size = Pt(18)
    addr_para.font.color.rgb = RGBColor(255, 255, 255)
    addr_para.alignment = PP_ALIGN.CENTER
    
    # 핵심 특징
    features = property_data.get('key_features', [])
    if features:
        feature_y = 6.0
        for feature in features[:3]:
            feat_box = slide.shapes.add_textbox(Inches(3.5), Inches(feature_y), Inches(3), Inches(0.3))
            feat_frame = feat_box.text_frame
            feat_frame.text = f"• {feature}"
            feat_para = feat_frame.paragraphs[0]
            feat_para.font.size = Pt(16)
            feat_para.font.color.rgb = RGBColor(255, 255, 255)
            feat_para.alignment = PP_ALIGN.CENTER
            feature_y += 0.35
    
    # 2. 건물 정보 슬라이드
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 배경색
    background2 = slide2.background
    fill2 = background2.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(245, 245, 245)
    
    # 제목
    title2_box = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.6))
    title2_frame = title2_box.text_frame
    title2_frame.text = "건물 정보"
    title2_para = title2_frame.paragraphs[0]
    title2_para.font.size = Pt(36)
    title2_para.font.bold = True
    title2_para.font.color.rgb = RGBColor(26, 35, 53)
    
    # 건축 개요 테이블 (좌측)
    left_x = 1
    info_y = 1.5
    line_height = 0.35
    
    building_info = property_data.get('building_info', {})
    
    info_items = [
        ("소재지", building_info.get('address', '')),
        ("지역지구", building_info.get('land_use_zone', '')),
        ("연면적", building_info.get('total_area', '')),
        ("대지면적", building_info.get('building_area', '')),
        ("건폐율", building_info.get('building_coverage', '')),
        ("용적률", building_info.get('floor_area_ratio', '')),
        ("건물규모", f"지상 {building_info.get('floors_above', '')}층 / 지하 {building_info.get('floors_below', '')}층"),
        ("주차대수", building_info.get('parking_spaces', '')),
        ("승강기", building_info.get('elevators', ''))
    ]
    
    for label, value in info_items:
        # 라벨
        label_box = slide2.shapes.add_textbox(Inches(left_x), Inches(info_y), Inches(2), Inches(line_height))
        label_frame = label_box.text_frame
        label_frame.text = label
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(14)
        label_para.font.bold = True
        label_para.font.color.rgb = RGBColor(60, 60, 60)
        
        # 값
        value_box = slide2.shapes.add_textbox(Inches(left_x + 2), Inches(info_y), Inches(3), Inches(line_height))
        value_frame = value_box.text_frame
        value_frame.text = str(value)
        value_para = value_frame.paragraphs[0]
        value_para.font.size = Pt(14)
        value_para.font.color.rgb = RGBColor(60, 60, 60)
        
        info_y += line_height
    
    # 3. 입지 분석 슬라이드
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    background3 = slide3.background
    fill3 = background3.fill
    fill3.solid()
    fill3.fore_color.rgb = RGBColor(245, 245, 245)
    
    title3_box = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.6))
    title3_frame = title3_box.text_frame
    title3_frame.text = "위치 분석"
    title3_para = title3_frame.paragraphs[0]
    title3_para.font.size = Pt(36)
    title3_para.font.bold = True
    title3_para.font.color.rgb = RGBColor(26, 35, 53)
    
    # 입지 설명
    location_desc = property_data.get('location_description', '')
    if location_desc:
        desc_box = slide3.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(3))
        desc_frame = desc_box.text_frame
        desc_frame.text = location_desc
        desc_frame.word_wrap = True
        desc_para = desc_frame.paragraphs[0]
        desc_para.font.size = Pt(18)
        desc_para.font.color.rgb = RGBColor(60, 60, 60)
        desc_para.line_spacing = 1.5
    
    # 저장
    prs.save(output_path)
    return output_path


def generate_print_pptx(property_data: dict, output_path: str):
    """
    인쇄용 PPTX 제안서 생성
    간결한 표 형식, 출력 최적화
    """
    # 템플릿 로드
    template_path = os.path.join(os.path.dirname(__file__), '..', 'template_print.pptx')
    
    if os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    
    # 기존 슬라이드 삭제
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    
    # 1. 표지 (인쇄용 - 심플)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 흰색 배경
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # 건물명
    title_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = property_data.get('building_name', '건물명')
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 0, 0)
    title_para.alignment = PP_ALIGN.CENTER
    
    # 매각금액
    price_box = slide.shapes.add_textbox(Inches(2), Inches(3.8), Inches(6), Inches(0.6))
    price_frame = price_box.text_frame
    price_frame.text = f"매각금액 : {property_data.get('sale_price', '')}원"
    price_para = price_frame.paragraphs[0]
    price_para.font.size = Pt(24)
    price_para.font.bold = True
    price_para.font.color.rgb = RGBColor(0, 0, 0)
    price_para.alignment = PP_ALIGN.CENTER
    
    # 회사 정보
    company_box = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(6), Inches(1))
    company_frame = company_box.text_frame
    company_info = property_data.get('company_info', '이음프로퍼티 부동산중개법인')
    company_frame.text = company_info
    company_para = company_frame.paragraphs[0]
    company_para.font.size = Pt(18)
    company_para.font.color.rgb = RGBColor(60, 60, 60)
    company_para.alignment = PP_ALIGN.CENTER
    
    # 2. 건축 개요 (표 형식)
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    background2 = slide2.background
    fill2 = background2.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(255, 255, 255)
    
    # 제목
    title2_box = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.5))
    title2_frame = title2_box.text_frame
    title2_frame.text = "건축 개요"
    title2_para = title2_frame.paragraphs[0]
    title2_para.font.size = Pt(28)
    title2_para.font.bold = True
    title2_para.font.color.rgb = RGBColor(0, 0, 0)
    
    # 표 추가
    rows = 12
    cols = 4
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(5)
    
    table = slide2.shapes.add_table(rows, cols, left, top, width, height).table
    
    # 컬럼 너비 설정
    table.columns[0].width = Inches(1.5)
    table.columns[1].width = Inches(3)
    table.columns[2].width = Inches(1)
    table.columns[3].width = Inches(2.5)
    
    building_info = property_data.get('building_info', {})
    
    # 표 데이터
    table_data = [
        ["소재지", building_info.get('address', ''), "", ""],
        ["지역지구", building_info.get('land_use_zone', ''), "", ""],
        ["주용도", building_info.get('main_purpose', ''), "", ""],
        ["대지면적", building_info.get('building_area', ''), "/", ""],
        ["연면적", building_info.get('total_area', ''), "/", ""],
        ["건폐율", building_info.get('building_coverage', ''), "", ""],
        ["용적률", building_info.get('floor_area_ratio', ''), "", ""],
        ["주요구조", building_info.get('structure', ''), "", ""],
        ["건물규모", f"지상 {building_info.get('floors_above', '')} / 지하 {building_info.get('floors_below', '')}", "", ""],
        ["주차대수", building_info.get('parking_spaces', ''), "", ""],
        ["승강기", building_info.get('elevators', ''), "", ""],
        ["준공일", building_info.get('completion_date', ''), "", ""]
    ]
    
    # 표 채우기
    for i, row_data in enumerate(table_data):
        for j, cell_value in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_value)
            
            # 셀 스타일링
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = RGBColor(0, 0, 0)
                
                # 첫 번째 컬럼은 볼드
                if j == 0:
                    paragraph.font.bold = True
    
    # 저장
    prs.save(output_path)
    return output_path
